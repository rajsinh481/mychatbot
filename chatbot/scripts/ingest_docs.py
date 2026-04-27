import os
from pathlib import Path
from docx import Document
import pandas as pd
from pptx import Presentation

# =========================
# PATHS
# =========================
ROOT_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = ROOT_DIR / "knowledge" / "docs"
OUTPUT_FILE = ROOT_DIR / "knowledge" / "extracted_docs_text.txt"


def read_docx(file_path: Path) -> str:
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def read_excel(file_path: Path) -> str:
    excel_data = pd.read_excel(file_path, sheet_name=None)
    all_text = []

    for sheet_name, df in excel_data.items():
        all_text.append(f"\n=== SHEET: {sheet_name} ===\n")
        all_text.append(df.fillna("").to_string(index=False))

    return "\n".join(all_text)


def read_ppt(file_path: Path) -> str:
    prs = Presentation(file_path)
    text = []

    for i, slide in enumerate(prs.slides, start=1):
        text.append(f"\n=== SLIDE {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())

    return "\n".join(text)


def process_files():
    if not DOCS_DIR.exists():
        print("❌ knowledge/docs folder not found.")
        return

    files = os.listdir(DOCS_DIR)
    all_text = ""

    supported = [".docx", ".xlsx", ".pptx"]
    matched_files = [f for f in files if any(f.lower().endswith(ext) for ext in supported)]

    if not matched_files:
        print("❌ No DOCX / PPTX / XLSX files found in knowledge/docs")
        return

    for file in matched_files:
        file_path = DOCS_DIR / file
        print(f"Processing: {file}")

        try:
            if file.lower().endswith(".docx"):
                extracted_text = read_docx(file_path)

            elif file.lower().endswith(".xlsx"):
                extracted_text = read_excel(file_path)

            elif file.lower().endswith(".pptx"):
                extracted_text = read_ppt(file_path)

            else:
                continue

            all_text += f"\n\n===== DOCUMENT SOURCE: {file} =====\n"
            all_text += extracted_text
            all_text += "\n"

        except Exception as e:
            print(f"❌ Error in {file}: {e}")
            continue

    OUTPUT_FILE.write_text(all_text, encoding="utf-8")
    print("✅ Docs text saved to extracted_docs_text.txt")


if __name__ == "__main__":
    process_files()